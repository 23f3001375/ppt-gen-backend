# core/generator.py

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

def create_ppt_from_template(slide_data, output_path, template_path=None, template_style=None):
    """
    Creates a PPT file from slide_data, applying styles from a template.
    This version is more robust for handling non-standard templates.
    """
    prs = Presentation(template_path) if template_path else Presentation()

    # --- Find a suitable "Title and Content" layout ---
    title_and_content_layout = None
    for layout in prs.slide_layouts:
        has_title = any(ph.placeholder_format.type == PP_PLACEHOLDER.TITLE for ph in layout.placeholders)
        # Check for Body OR a generic Object placeholder, which is common
        has_body = any(ph.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT) for ph in layout.placeholders)
        if has_title and has_body:
            title_and_content_layout = layout
            break

    if not title_and_content_layout:
        title_and_content_layout = prs.slide_layouts[1]

    # --- Create slides ---
    for item in slide_data:
        slide = prs.slides.add_slide(title_and_content_layout)
        
        title_shape = slide.shapes.title
        body_shape = None

        # --- NEW, MORE ROBUST LOGIC TO FIND THE BODY SHAPE ---
        # Find the title placeholder
        if title_shape:
            title_shape.text = item.get("title", "No Title")

        # Iterate through all placeholders on the slide to find the body
        for shape in slide.placeholders:
            # The title placeholder is handled, so we look for the next available one
            if shape.placeholder_format.type != PP_PLACEHOLDER.TITLE:
                body_shape = shape
                break
        
        # --- End of new logic ---
        
        if body_shape:
            tf = body_shape.text_frame
            tf.clear()
            
            points = item.get("points", [])
            if points:
                p = tf.paragraphs[0]
                p.text = points[0]
                
                for point_text in points[1:]:
                    p = tf.add_paragraph()
                    p.text = point_text
                    p.level = 0

    prs.save(output_path)
    return output_path